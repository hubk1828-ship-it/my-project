"""Rebuild PPTX v2 - handles groups, images, tables, text"""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os, traceback

desktop = os.path.join(os.environ['USERPROFILE'], 'Desktop')
source = os.path.join(desktop, 'REPAIRED_PYTHON.pptx')
output = os.path.join(desktop, 'REBUILT_CLEAN_V2.pptx')
temp_img = os.path.join(desktop, 'pptx_test', 'temp_img')

print("Opening source...")
src_prs = Presentation(source)

new_prs = Presentation()
new_prs.slide_width = src_prs.slide_width
new_prs.slide_height = src_prs.slide_height
blank_layout = new_prs.slide_layouts[6]

img_counter = 0

def add_image_to_slide(slide, shape):
    global img_counter
    try:
        image = shape.image
        ct = image.content_type
        if 'svg' in ct:
            return False
        ext = ct.split('/')[-1].replace('jpeg','jpg')
        img_counter += 1
        path = f"{temp_img}_{img_counter}.{ext}"
        with open(path, 'wb') as f:
            f.write(image.blob)
        slide.shapes.add_picture(path, shape.left, shape.top, shape.width, shape.height)
        os.remove(path)
        return True
    except:
        return False

def copy_text_shape(slide, shape):
    try:
        txBox = slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        for p_idx, src_para in enumerate(shape.text_frame.paragraphs):
            para = tf.add_paragraph() if p_idx > 0 else tf.paragraphs[0]
            if src_para.alignment:
                para.alignment = src_para.alignment
            
            for run in src_para.runs:
                new_run = para.add_run()
                new_run.text = run.text
                try:
                    if run.font.size: new_run.font.size = run.font.size
                    if run.font.bold is not None: new_run.font.bold = run.font.bold
                    if run.font.italic is not None: new_run.font.italic = run.font.italic
                    if run.font.name: new_run.font.name = run.font.name
                    try:
                        if run.font.color and run.font.color.rgb:
                            new_run.font.color.rgb = run.font.color.rgb
                    except: pass
                except: pass
        return True
    except:
        return False

def copy_table(slide, shape):
    try:
        table = shape.table
        rows, cols = len(table.rows), len(table.columns)
        new_shape = slide.shapes.add_table(rows, cols, shape.left, shape.top, shape.width, shape.height)
        new_table = new_shape.table
        for r in range(rows):
            for c in range(cols):
                try:
                    new_table.cell(r, c).text = table.cell(r, c).text
                except: pass
        return True
    except:
        return False

def process_shape(slide, shape):
    """Process a single shape and add it to the slide"""
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return add_image_to_slide(slide, shape)
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            # Process group shapes recursively
            count = 0
            for child in shape.shapes:
                if process_shape(slide, child):
                    count += 1
            return count > 0
        elif shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:  # Only add non-empty text
                return copy_text_shape(slide, shape)
            return False
        elif shape.has_table:
            return copy_table(slide, shape)
        return False
    except:
        return False

print("Rebuilding slides...")
stats = {'images': 0, 'text': 0, 'tables': 0, 'groups': 0, 'failed': 0}

for slide_idx, src_slide in enumerate(src_prs.slides):
    slide_num = slide_idx + 1
    new_slide = new_prs.slides.add_slide(blank_layout)
    
    slide_shapes = 0
    for shape in src_slide.shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                if add_image_to_slide(new_slide, shape):
                    stats['images'] += 1
                    slide_shapes += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for child in shape.shapes:
                    if process_shape(new_slide, child):
                        slide_shapes += 1
                stats['groups'] += 1
            elif shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    if copy_text_shape(new_slide, shape):
                        stats['text'] += 1
                        slide_shapes += 1
            elif shape.has_table:
                if copy_table(new_slide, shape):
                    stats['tables'] += 1
                    slide_shapes += 1
        except Exception as e:
            stats['failed'] += 1
    
    if slide_num <= 3 or slide_num % 10 == 0 or slide_num == 63:
        print(f"  Slide {slide_num}: {slide_shapes} elements added")

print(f"\nStats: {stats}")
print(f"Saving...")
new_prs.save(output)
size_mb = os.path.getsize(output) / (1024*1024)
print(f"Saved REBUILT_CLEAN_V2.pptx: {size_mb:.2f} MB")
